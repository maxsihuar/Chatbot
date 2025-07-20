from langchain_community.document_loaders.text import TextLoader
from langchain_community.document_loaders import PyPDFLoader
from langchain_community.document_loaders import CSVLoader
from pptx import Presentation
from langchain_core.documents import Document
from langchain_community. document_loaders.base import BaseLoader

from qdrant_client import QdrantClient
from qdrant_client.models import VectorParams, Distance
from langchain_community.vectorstores import Qdrant
import google.generativeai as genai
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from dotenv import load_dotenv

from langchain.text_splitter import RecursiveCharacterTextSplitter

from collections import deque


import os
import sys

def LoadPPTX(pathArchivo: str) ->list[Document]| None:
    """
    Extrae texto de todas las diapositivas de un archivo .pptx

    Args:
        pathArchivo (str): Ruta del archivo .pptx

    Returns:
        list[str]: Lista de textos extraídos por cada shape
    """
    pptx = Presentation(pathArchivo)
    datos = []
    for slide in pptx.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texto = shape.text.strip()
                if texto:
                    datos.append(texto)

    documento = [
        Document(
            page_content="\n".join(datos),
            metadata={"source":pathArchivo}
        )
    ]
            

    return documento

def LoadAll() -> list | None:
    """
    Crea una lista de documentos los cuales se abstraen de la carpeta Doc
    Lee archivos con extensiones .md, .txt, .pdf y .csv, y los convierte en objetos
    Document compatibles con LangChain para procesamiento posterior.

    Args:
        None
    
    Returns:
        archivos(list()): Retorna una lista de archivos Laingchain
    """
    path = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "Doc"))

    data = os.listdir(path)

    archivos = list()

    for archivo in data:
        pathArchivo = os.path.join(path,archivo)
        print(pathArchivo)
        
        if str.endswith(pathArchivo, ".md") or str.endswith(pathArchivo, ".txt"):
            documentos = TextLoader(pathArchivo, encoding="utf-8").load()
        
        elif str.endswith(pathArchivo, ".pdf"):
            documentos = PyPDFLoader(pathArchivo).load()
        elif str.endswith(pathArchivo, ".csv"):
            documentos = CSVLoader(pathArchivo, encoding="utf-8").load()
        elif str.endswith(pathArchivo,".pptx"):
            documentos = LoadPPTX(pathArchivo)
        else:
            continue

        archivos.extend(documentos)

    return archivos